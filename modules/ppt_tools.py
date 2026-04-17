import copy
from pathlib import Path
from pptx import Presentation

from utils.logger import log_action, log_failure
from utils.ui import print_warning


def _best_blank_layout(prs: Presentation):
    """
    Find safest blank layout.
    """

    for layout in prs.slide_layouts:
        try:
            if "blank" in layout.name.lower():
                return layout
        except Exception:
            continue

    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]

    return prs.slide_layouts[0]


def _remove_default_placeholders(slide) -> None:
    """
    Remove placeholder objects from newly created slide.
    """

    try:
        for shape in list(slide.shapes):
            try:
                if getattr(shape, "is_placeholder", False):
                    el = shape.element
                    el.getparent().remove(el)
            except Exception:
                continue
    except Exception:
        pass


def merge_pptx(
    file_list: list[str],
    output_path: str
) -> None:
    """
    Merge multiple PPTX files into one presentation.

    Uses first file as base template/theme.

    Best for:
    - text slides
    - image slides
    - normal shapes

    Limitations may apply to:
    - animations
    - transitions
    - SmartArt
    - charts
    - media
    - macros
    """

    if not file_list:
        raise ValueError("No PPTX files provided.")

    files_used = 0
    files_skipped = 0

    base_slides = 0
    imported_slides = 0
    total_shapes_copied = 0
    total_shapes_skipped = 0

    # ---------------------------------------------------------
    # Load base presentation
    # ---------------------------------------------------------
    try:
        base_prs = Presentation(file_list[0])
        files_used += 1
        base_slides = len(base_prs.slides)

    except Exception as e:
        raise RuntimeError(
            f"Failed to open base PPTX: {file_list[0]} ({e})"
        )

    blank_layout = _best_blank_layout(base_prs)

    # ---------------------------------------------------------
    # Merge remaining files
    # ---------------------------------------------------------
    for file_path in file_list[1:]:
        try:
            prs = Presentation(file_path)

            if len(prs.slides) == 0:
                msg = f"Empty PPTX skipped: {file_path}"
                print_warning(msg)
                log_failure(msg)
                files_skipped += 1
                continue

            for slide in prs.slides:
                try:
                    new_slide = base_prs.slides.add_slide(blank_layout)

                    _remove_default_placeholders(new_slide)

                    for shape in slide.shapes:
                        try:
                            new_element = copy.deepcopy(shape.element)

                            new_slide.shapes._spTree.insert_element_before(
                                new_element,
                                "p:extLst"
                            )

                            total_shapes_copied += 1

                        except Exception:
                            total_shapes_skipped += 1
                            continue

                    imported_slides += 1

                except Exception:
                    continue

            files_used += 1

        except Exception as e:
            msg = f"PPTX skipped: {file_path} ({e})"
            print_warning(msg)
            log_failure(msg)
            files_skipped += 1

    # ---------------------------------------------------------
    # Validate final result
    # ---------------------------------------------------------
    total_slides = len(base_prs.slides)

    if total_slides == 0:
        raise ValueError("No valid slides available to save.")

    # ---------------------------------------------------------
    # Save
    # ---------------------------------------------------------
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    try:
        base_prs.save(output_path)

    except Exception as e:
        raise RuntimeError(f"Failed to save merged PPTX: {e}")

    # ---------------------------------------------------------
    # Logging
    # ---------------------------------------------------------
    log_action(
        "PPTX merge success | "
        f"files_used={files_used} "
        f"files_skipped={files_skipped} "
        f"base_slides={base_slides} "
        f"imported_slides={imported_slides} "
        f"total_slides={total_slides} "
        f"shapes_copied={total_shapes_copied} "
        f"shapes_skipped={total_shapes_skipped} "
        f"-> {output_path}"
    )